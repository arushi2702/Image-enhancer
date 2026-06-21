from pathlib import Path

from docx import Document
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
DOCS_DIR = ROOT / "docs"
REPORT_MD = DOCS_DIR / "PROJECT_REPORT.md"
REPORT_DOCX = DOCS_DIR / "Image_Quality_Assessment_Report.docx"
PRESENTATION_PPTX = DOCS_DIR / "Image_Quality_Assessment_Presentation.pptx"


def build_docx():
    document = Document()
    markdown = REPORT_MD.read_text(encoding="utf-8").splitlines()

    for line in markdown:
        stripped = line.strip()

        if not stripped:
            document.add_paragraph("")
            continue

        if stripped.startswith("# "):
            paragraph = document.add_paragraph()
            run = paragraph.add_run(stripped[2:])
            run.bold = True
            run.font.size = Pt(20)
            paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            continue

        if stripped.startswith("## "):
            document.add_heading(stripped[3:], level=1)
            continue

        if stripped.startswith("### "):
            document.add_heading(stripped[4:], level=2)
            continue

        if stripped.startswith("- "):
            document.add_paragraph(stripped[2:], style="List Bullet")
            continue

        if stripped[:2].isdigit() and stripped[2:4] == ". ":
            document.add_paragraph(stripped[4:], style="List Number")
            continue

        document.add_paragraph(stripped)

    document.save(REPORT_DOCX)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color)


def add_title(slide, text, x, y, w, h, size=24, color="203047"):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string(color)
    p.font.name = "Aptos"


def add_body_lines(slide, lines, x, y, w, h, size=18, color="203047"):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()

    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = RGBColor.from_string(color)
        paragraph.font.name = "Aptos"
        paragraph.space_after = Pt(10)
        paragraph.bullet = True


def add_card(slide, title, body, x, y, w, h):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string("EEF4F8")
    shape.line.color.rgb = RGBColor.from_string("C5D4E3")

    title_box = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.1), w - Inches(0.2), Inches(0.24))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = title
    title_p.font.bold = True
    title_p.font.size = Pt(16)
    title_p.font.color.rgb = RGBColor.from_string("203047")
    title_p.font.name = "Aptos"

    body_box = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.4), w - Inches(0.2), h - Inches(0.45))
    body_p = body_box.text_frame.paragraphs[0]
    body_p.text = body
    body_p.font.size = Pt(11)
    body_p.font.color.rgb = RGBColor.from_string("5F7693")
    body_p.font.name = "Aptos"


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, "EAF1F8")
    add_title(
        slide,
        "Image Quality Assessment & Enhancement Web App",
        Inches(0.8),
        Inches(0.9),
        Inches(10.5),
        Inches(0.7),
        size=24,
    )
    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(8.5), Inches(1.2))
    para = subtitle.text_frame.paragraphs[0]
    para.text = "Assessment-first image enhancement using Flask, OpenCV, and Next.js"
    para.font.size = Pt(18)
    para.font.color.rgb = RGBColor.from_string("456B94")
    para.font.name = "Aptos"

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, "F7FAFD")
    add_title(slide, "Problem Statement", Inches(0.6), Inches(0.4), Inches(8), Inches(0.5))
    add_body_lines(
        slide,
        [
            "Many images suffer from low brightness, weak contrast, visible noise, and soft details.",
            "Applying the same enhancement to every image often over-processes already good inputs.",
            "This project solves that by assessing image quality first and enhancing only when necessary.",
        ],
        Inches(0.85),
        Inches(1.2),
        Inches(11.2),
        Inches(4.5),
    )

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, "F7FAFD")
    add_title(slide, "System Architecture", Inches(0.6), Inches(0.4), Inches(8), Inches(0.5))
    add_card(slide, "Frontend", "Next.js UI\nReact + TypeScript\nUpload, preview, compare, download", Inches(0.9), Inches(1.5), Inches(3.2), Inches(2.3))
    add_card(slide, "Backend", "Flask API\nOpenCV enhancement pipeline\nQuality scoring and decision logic", Inches(5.0), Inches(1.5), Inches(3.8), Inches(2.3))
    arrow = slide.shapes.add_textbox(Inches(4.2), Inches(2.2), Inches(0.5), Inches(0.5))
    ap = arrow.text_frame.paragraphs[0]
    ap.text = "→"
    ap.alignment = PP_ALIGN.CENTER
    ap.font.size = Pt(28)
    ap.font.bold = True
    ap.font.color.rgb = RGBColor.from_string("7F9CBB")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, "F7FAFD")
    add_title(slide, "Metrics Used for Assessment", Inches(0.6), Inches(0.4), Inches(8), Inches(0.5))
    cards = [
        ("Brightness", "Measures overall exposure"),
        ("Contrast", "Measures spread of grayscale values"),
        ("Sharpness", "Uses Laplacian variance for edge clarity"),
        ("Noise", "Estimated from blur residual"),
        ("Dynamic Range", "95th minus 5th percentile intensity"),
        ("Colorfulness", "Measures vividness of color distribution"),
    ]
    for idx, (title, body) in enumerate(cards):
        row = idx // 2
        col = idx % 2
        add_card(
            slide,
            title,
            body,
            Inches(0.8 + col * 4.35),
            Inches(1.3 + row * 1.55),
            Inches(3.75),
            Inches(1.0),
        )

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, "F7FAFD")
    add_title(slide, "OpenCV Enhancement Pipeline", Inches(0.6), Inches(0.4), Inches(8), Inches(0.5))
    add_body_lines(
        slide,
        [
            "Gray-world white balance to reduce global color cast.",
            "Adaptive gamma correction for exposure control.",
            "CLAHE on the luminance channel for local contrast enhancement.",
            "Selective luminance lift to brighten darker regions more than highlights.",
            "Adaptive denoising based on estimated noise level.",
            "Unsharp masking, highlight/shadow protection, and color preservation.",
        ],
        Inches(0.85),
        Inches(1.2),
        Inches(11.0),
        Inches(4.8),
    )

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, "F7FAFD")
    add_title(slide, "Assessment-First Decision Logic", Inches(0.6), Inches(0.4), Inches(8.5), Inches(0.5))
    add_body_lines(
        slide,
        [
            "The backend computes an overall quality score from multiple metrics.",
            "If the score is at least 90%, the image is preserved and no enhancement is applied.",
            "If the score is below 90%, the backend detects issues such as low brightness, soft details, or visible noise.",
            "Enhancement modes are selected according to detected issues instead of using one fixed formula.",
        ],
        Inches(0.85),
        Inches(1.2),
        Inches(11.0),
        Inches(4.8),
    )

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, "F7FAFD")
    add_title(slide, "User-Facing Output", Inches(0.6), Inches(0.4), Inches(8), Inches(0.5))
    add_body_lines(
        slide,
        [
            "Original and processed images are shown side by side.",
            "Overall quality score and rating are shown for both versions.",
            "Assessment panel explains whether enhancement was needed.",
            "Improvement section summarizes metric changes.",
            "The processed image can be downloaded directly from the interface.",
        ],
        Inches(0.85),
        Inches(1.2),
        Inches(11.0),
        Inches(4.8),
    )

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, "F7FAFD")
    add_title(slide, "Limitations and Future Scope", Inches(0.6), Inches(0.4), Inches(8.5), Inches(0.5))
    add_body_lines(
        slide,
        [
            "This is a classical computer vision pipeline, not a deep learning model.",
            "Severe blur or missing detail cannot be perfectly reconstructed.",
            "Future scope includes face-aware enhancement, document mode, object detection, and segmentation-assisted enhancement.",
        ],
        Inches(0.85),
        Inches(1.2),
        Inches(11.0),
        Inches(4.8),
    )

    prs.save(PRESENTATION_PPTX)


def main():
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    build_docx()
    build_presentation()
    print(REPORT_DOCX)
    print(PRESENTATION_PPTX)


if __name__ == "__main__":
    main()
